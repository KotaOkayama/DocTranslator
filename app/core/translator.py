#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPTX/DOCX/PDF Text Translation Module

This module extracts text from PPTX/DOCX/PDF files, translates it using the GenAI Hub API, and re-imports the translated text into new files.
It maintains formatting and hyperlinks as much as possible.
"""

import json
import os
import time
import logging
import re
import subprocess
import tempfile
import shutil
from typing import Optional, Tuple, Callable, Dict, List

import PyPDF2
import pdfplumber
from docx import Document
from docx.shared import Inches
import pandas as pd

# dotenv import
from dotenv import load_dotenv

import pandas as pd
import requests

import concurrent.futures

# Logging configuration
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler("logs/translator.log", encoding="utf-8"),
    ],
)
logger = logging.getLogger(__name__)

# Load .env file
load_dotenv()

# GenAI Hub API URL
def get_api_url():
    """Get API URL from environment variable or use default"""
    from app.config import get_api_url
    return get_api_url()

def get_api_key():
    """Get API key from environment variable"""
    from app.config import get_api_key
    return get_api_key()

# Keywords to filter out models not suitable for translation
EXCLUDED_MODEL_KEYWORDS = [
    "copilot",
    "documentation", 
    "jira",
    "ngsiem",
    "dom-data-extractor",
    "wiki",
    "compintel",
    "coding"
]

def filter_translation_models(models: List[str]) -> List[str]:
    """
    Filter models suitable for translation
    
    Args:
        models: List of model names
        
    Returns:
        List of model names suitable for translation
    """
    filtered_models = []
    
    for model in models:
        model_lower = model.lower()
        
        # Check if excluded keywords are present
        should_exclude = any(keyword in model_lower for keyword in EXCLUDED_MODEL_KEYWORDS)
        
        if not should_exclude:
            filtered_models.append(model)
    
    logger.info(f"Model filtering result: {len(models)} -> {len(filtered_models)} models")
    logger.debug(f"Excluded models: {set(models) - set(filtered_models)}")
    
    return filtered_models

def format_model_name(model_name: str) -> str:
    """
    Format GenAI Hub specific model names appropriately
    
    Args:
        model_name: Original model name
        
    Returns:
        Formatted model name
    """
    # Format Claude models
    if model_name.startswith("claude-"):
        # claude-3-5-sonnet-v2 -> Claude 3.5 Sonnet V2
        # claude-4-sonnet -> Claude 4 Sonnet
        # claude-3-7-sonnet -> Claude 3.7 Sonnet
        parts = model_name.split("-")
        if len(parts) >= 3:
            version = parts[1]  # 3, 4
            subversion = parts[2] if len(parts) > 2 else ""  # 5, 7
            model_type = parts[3] if len(parts) > 3 else parts[2]  # sonnet, haiku, opus
            variant = parts[4] if len(parts) > 4 else ""  # v2
            
            formatted = f"Claude {version}"
            if subversion and subversion.isdigit():
                formatted += f".{subversion}"
            formatted += f" {model_type.title()}"
            if variant:
                formatted += f" {variant.upper()}"
            
            return formatted
    
    # Format Llama models
    elif model_name.startswith("llama"):
        # llama3-1-70b -> Llama 3.1 70B
        # llama4-maverick-17b -> Llama 4 Maverick 17B
        parts = model_name.split("-")
        if len(parts) >= 2:
            version_part = parts[0].replace("llama", "")  # 3, 4
            
            formatted = f"Llama {version_part}"
            
            # Process version number
            if len(parts) > 1 and parts[1].isdigit():
                formatted += f".{parts[1]}"
                remaining_parts = parts[2:]
            else:
                remaining_parts = parts[1:]
            
            # Process remaining parts
            for part in remaining_parts:
                if part.endswith("b") and part[:-1].isdigit():
                    # Size information (70b -> 70B)
                    formatted += f" {part.upper()}"
                else:
                    # Other information (maverick, scout, etc.)
                    formatted += f" {part.title()}"
            
            return formatted
    
    # Capitalize first letter for other models
    return model_name.replace("-", " ").title()

def fetch_available_models() -> Dict[str, str]:
    """
    Fetch available models from GenAI HUB
    
    Returns:
        Model dictionary {model_id: display_name}
        
    Raises:
        ValueError: If API call fails
        ConnectionError: If connection fails
    """
    logger.info("Fetching models from GenAI HUB...")
    
    # Get API settings
    api_key = get_api_key()
    api_url = get_api_url()
    
    if not api_key:
        raise ValueError("API key is not set. Please configure API key in settings.")
    
    if not api_url:
        raise ValueError("API URL is not set. Please configure API URL in settings.")
    
    # Build models endpoint URL
    # Convert chat endpoint to models endpoint
    if api_url.endswith("/chat/completions"):
        models_url = api_url.replace("/chat/completions", "/models")
    elif api_url.endswith("/v1/chat/completions"):
        models_url = api_url.replace("/v1/chat/completions", "/v1/models")
    else:
        # Fallback: append /models to URL
        models_url = api_url.rstrip("/") + "/models"
    
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    
    try:
        logger.debug(f"Models fetch URL: {models_url}")
        response = requests.get(
            models_url,
            headers=headers,
            timeout=30
        )
        
        logger.debug(f"API response status: {response.status_code}")
        
        if response.status_code != 200:
            logger.error(f"Models fetch error: status code {response.status_code}")
            logger.error(f"Response: {response.text}")
            raise ValueError(f"Failed to fetch models. Status code: {response.status_code}")
        
        result = response.json()
        logger.debug(f"API response: {result}")
        
        # Determine response format and parse appropriately
        models_list = None
        
        # New format (OpenAI compatible): uses data field
        if "data" in result and isinstance(result["data"], list):
            logger.debug("Detected OpenAI compatible response format")
            models_list = [model["id"] for model in result["data"] if "id" in model]
        
        # Old format: uses content field (for backward compatibility)
        elif "content" in result and isinstance(result["content"], list):
            logger.debug("Detected old response format")
            models_list = result["content"]
        
        # Unknown format
        else:
            logger.error(f"Unknown response format: {result}")
            raise ValueError("Failed to fetch models. Invalid response format.")
        
        if not isinstance(models_list, list):
            logger.error(f"Models list is not a list: {type(models_list)}")
            raise ValueError("Failed to fetch models. Invalid response format.")
        
        logger.info(f"Fetched models count: {len(models_list)}")
        logger.debug(f"Fetched models: {models_list}")
        
        # Filter models suitable for translation
        filtered_models = filter_translation_models(models_list)
        
        if not filtered_models:
            logger.warning("No models suitable for translation found")
            raise ValueError("No models suitable for translation found.")
        
        # Create model dictionary and sort alphabetically
        models_dict = {}
        for model in filtered_models:
            display_name = format_model_name(model)
            models_dict[model] = display_name
        
        # Sort by model ID alphabetically
        sorted_models_dict = dict(sorted(models_dict.items()))
        
        logger.info(f"Available translation models (alphabetical order): {list(sorted_models_dict.keys())}")
        
        return sorted_models_dict
        
    except requests.exceptions.RequestException as e:
        logger.error(f"Models fetch request error: {e}")
        raise ConnectionError(f"Failed to connect to GenAI HUB: {str(e)}")
    except json.JSONDecodeError as e:
        logger.error(f"JSON decode error: {e}")
        logger.error(f"Response: {response.text if 'response' in locals() else 'None'}")
        raise ValueError(f"Failed to fetch models. Response parse error: {str(e)}")
    except Exception as e:
        logger.error(f"Models fetch error: {e}")
        raise ValueError(f"Failed to fetch models: {str(e)}")

def get_available_models() -> Dict[str, str]:
    """
    Get available models (no fallback)
    
    Returns:
        Model dictionary {model_id: display_name}
    Raises:
        ValueError: If API call fails
        ConnectionError: If connection fails
    """
    return fetch_available_models()

# Language options
LANGUAGES = {
    "ja": "Japanese",
    "en": "English",
    "zh": "Chinese",
    "ko": "Korean",
    "fr": "French",
    "de": "German",
    "es": "Spanish",
    "hi": "Hindi",
    "vi": "Vietnamese",
    "th": "Thai",
}

# Translation Settings
TRANSLATION_CONFIG = {
    "chunk_size": 1500,  # Maximum length of text to split
    "batch_size": 5,  # Maximum number of texts to batch process
    "parallel_workers": 4,  # Maximum number of parallel processes
    "use_cache": True,  # Whether to use cache
    "optimize_text": True,  # Whether to optimize text
    "api_timeout": 30,  # API call timeout (seconds)
    "retry_count": 2,  # Number of retries on error
}

# Global Translation Cache
TRANSLATION_CACHE = {}

def save_translation_cache(cache_file: str = "translation_cache.json"):
    """Save translation cache to file"""
    try:
        with open(cache_file, "w", encoding="utf-8") as f:
            json.dump(TRANSLATION_CACHE, f, ensure_ascii=False, indent=2)
        logger.info(f"Translation cache saved: {len(TRANSLATION_CACHE)} entries")
    except Exception as e:
        logger.error(f"Failed to save cache: {e}")

def load_translation_cache(cache_file: str = "translation_cache.json"):
    """Load translation cache from file"""
    global TRANSLATION_CACHE
    try:
        if os.path.exists(cache_file):
            with open(cache_file, "r", encoding="utf-8") as f:
                TRANSLATION_CACHE = json.load(f)
            logger.info(f"Translation cache loaded: {len(TRANSLATION_CACHE)} entries")
    except Exception as e:
        logger.error(f"Failed to load cache: {e}")
        TRANSLATION_CACHE = {}

# Load cache
load_translation_cache()

def translate_text_with_genai_hub(
    text: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    ai_instruction: str = "",
) -> str:
    """
    Translate text using GenAI Hub API.
    """
    if not text.strip():
        return ""
    
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()
        logger.debug(f"Using default model: {model}")

    # Get API key and URL from environment variables
    if not api_key:
        api_key = os.environ.get("GENAI_HUB_API_KEY")

    # Get API URL from environment variable
    api_url = os.environ.get("GENAI_HUB_API_URL")

    if not api_key:
        raise ValueError(
            "API key is not set. Please check environment variable GENAI_HUB_API_KEY."
        )
    
    if not api_url:
        raise ValueError(
            "API URL is not set. Please check environment variable GENAI_HUB_API_URL."
        )

    # Log only first and last few characters of API key (for security)
    masked_api_key = f"{api_key[:4]}...{api_key[-4:]}" if len(api_key) > 8 else "****"
    logger.info(
        f"Translation: {source_lang} -> {target_lang}, Model: {model}, API Key: {masked_api_key}, API URL: {api_url}"
    )
    logger.debug(f"Translation text: {text[:100]}...")

    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}

    # Create prompt with translation instructions
    prompt = f"""Translate the following text from {LANGUAGES.get(source_lang, source_lang)} to {LANGUAGES.get(target_lang, target_lang)}.
Maintain the original text format and only perform translation.
No explanations or annotations other than translation are needed.
Preserve formatting such as line breaks and spaces as much as possible.
"""

    # Add supplementary instructions if provided
    if ai_instruction:
        prompt += f"\nSupplementary instructions: {ai_instruction}\n"

    prompt += f"""
Text:
{text}

Translation:"""

    payload = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.1,  # Set low temperature for consistent translation
        "max_tokens": 4000,
    }

    # Retry counter
    retry_count = 0
    max_retries = TRANSLATION_CONFIG["retry_count"]

    while retry_count <= max_retries:
        try:
            logger.debug(f"Sending API request: {api_url}")
            response = requests.post(
                api_url,
                headers=headers,
                json=payload,
                timeout=TRANSLATION_CONFIG["api_timeout"],
            )

            # Log response details
            logger.debug(f"API response status: {response.status_code}")

            # Error check
            if response.status_code != 200:
                logger.error(f"API error: status code {response.status_code}")
                logger.error(f"Response: {response.text}")

                if retry_count < max_retries:
                    retry_count += 1
                    wait_time = 2**retry_count  # Exponential backoff
                    logger.info(
                        f"Retry {retry_count}/{max_retries} in {wait_time} seconds"
                    )
                    time.sleep(wait_time)
                    continue

                return f"[Translation error: API status code {response.status_code}]"

            response.raise_for_status()

            result = response.json()
            if "choices" not in result or len(result["choices"]) == 0:
                logger.error(f"API response missing 'choices': {result}")

                if retry_count < max_retries:
                    retry_count += 1
                    wait_time = 2**retry_count
                    logger.info(
                        f"Retry {retry_count}/{max_retries} in {wait_time} seconds"
                    )
                    time.sleep(wait_time)
                    continue

                return "[Translation error: Invalid response]"

            translated_text = result["choices"][0]["message"]["content"].strip()

            # Remove prefix like "Translation:" if present
            if "Translation:" in translated_text:
                translated_text = translated_text.split("Translation:", 1)[1].strip()

            logger.debug(f"Translation result: {translated_text[:100]}...")
            return translated_text

        except requests.exceptions.RequestException as e:
            logger.error(f"API request error: {e}")

            if retry_count < max_retries:
                retry_count += 1
                wait_time = 2**retry_count
                logger.info(
                    f"Retry {retry_count}/{max_retries} in {wait_time} seconds"
                )
                time.sleep(wait_time)
                continue

            return f"[Translation error: Request failed - {str(e)}]"

        except json.JSONDecodeError as e:
            logger.error(f"JSON decode error: {e}")
            logger.error(
                f"Response: {response.text if 'response' in locals() else 'None'}"
            )

            if retry_count < max_retries:
                retry_count += 1
                wait_time = 2**retry_count
                logger.info(
                    f"Retry {retry_count}/{max_retries} in {wait_time} seconds"
                )
                time.sleep(wait_time)
                continue

            return "[Translation error: Invalid JSON response]"
        except Exception as e:
            logger.error(f"Translation error: {e}")
            logger.error(
                f"Response: {response.text if 'response' in locals() else 'None'}"
            )

            if retry_count < max_retries:
                retry_count += 1
                wait_time = 2**retry_count
                logger.info(
                    f"Retry {retry_count}/{max_retries} in {wait_time} seconds"
                )
                time.sleep(wait_time)
                continue

            return f"[Translation error: {str(e)}]"

def translate_text_with_cache(
    text: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    ai_instruction: str = "",
) -> str:
    """
    Translate text using cache.

    Args:
        text: Text to translate
        api_key: GenAI Hub API key
        model: Model name to use
        source_lang: Source language code
        target_lang: Target language code
        ai_instruction: Supplementary instructions for AI

    Returns:
        Translated text
    """
    if not text.strip():
        return ""
    
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()

    # Create cache key (including supplementary instructions)
    cache_key = f"{text}|{model}|{source_lang}|{target_lang}|{ai_instruction}"

    # Check cache
    if cache_key in TRANSLATION_CACHE:
        logger.debug(f"Retrieved translation from cache: {text[:30]}...")
        return TRANSLATION_CACHE[cache_key]

    # Translate if not in cache
    translated = translate_text_with_genai_hub(
        text, api_key, model, source_lang, target_lang, ai_instruction
    )

    # Save to cache
    TRANSLATION_CACHE[cache_key] = translated

    return translated

def optimize_text_for_translation(text: str) -> str:
    """
    Optimize text for translation.

    Args:
        text: Text to optimize

    Returns:
        Optimized text
    """
    # Remove blank lines
    lines = [line for line in text.split("\n") if line.strip()]

    # Skip lines with only numbers or symbols
    filtered_lines = []
    for line in lines:
        # Check if text needs translation (contains alphabets, kanji, hiragana, katakana, etc.)
        if any(c.isalpha() for c in line):
            filtered_lines.append(line)
        else:
            # Add lines with only numbers or symbols as is
            filtered_lines.append(line)

    return "\n".join(filtered_lines)

def translate_dataframe(
    df: pd.DataFrame,
    api_key: str,
    model: str,
    source_lang: str = "en",
    target_lang: str = "ja",
    progress_callback: Optional[Callable] = None,
    ai_instruction: str = "",
    check_cancelled: Optional[Callable] = None,
) -> pd.DataFrame:
    """
    Translate text in DataFrame in parallel.

    Args:
        df: DataFrame containing text to translate
        api_key: GenAI Hub API key
        model: Model name to use
        source_lang: Source language code
        target_lang: Target language code
        progress_callback: Callback function to report progress
        ai_instruction: Supplementary instructions for AI
        check_cancelled: Function to check cancellation status

    Returns:
        DataFrame containing translated text
    """
    logger.info(f"Translating text in DataFrame in parallel: {len(df)} items")

    # Optimize text
    if TRANSLATION_CONFIG["optimize_text"]:
        optimized_texts = [
            optimize_text_for_translation(text) for text in df["original_text"]
        ]
    else:
        optimized_texts = df["original_text"].tolist()

    # Parallel processing settings
    max_workers = TRANSLATION_CONFIG.get("parallel_workers", 4)
    batch_size = TRANSLATION_CONFIG.get("batch_size", 5)

    # Counter to track progress
    total_items = len(df)
    completed_items = 0
    translations = [""] * total_items

    # Thread-safe progress update function
    def update_progress(count=1):
        nonlocal completed_items
        completed_items += count
        if progress_callback:
            try:
                progress = 0.3 + (completed_items / total_items * 0.5)
                progress_callback(min(progress, 0.8))
            except Exception as e:
                logger.warning(f"Progress callback error: {e}")

    # Wrapper function for single text translation
    def translate_single_text(index, text):
        try:
            # Check cancellation
            if check_cancelled and check_cancelled():
                logger.info(f"Translation cancelled (item {index+1})")
                return index, "[Translation cancelled]"

            translated_text = translate_text_with_cache(
                text, api_key, model, source_lang, target_lang, ai_instruction
            )
            return index, translated_text
        except Exception as e:
            logger.error(f"Text translation error (item {index+1}): {e}")
            return index, f"[Translation error: {str(e)}]"

    # Parallel translation using batch processing
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Process by batch
        for batch_start in range(0, total_items, batch_size):
            # Check cancellation
            if check_cancelled and check_cancelled():
                logger.info("Translation cancelled")
                break

            batch_end = min(batch_start + batch_size, total_items)
            batch_texts = optimized_texts[batch_start:batch_end]
            batch_indices = list(range(batch_start, batch_end))

            # Submit batch translation tasks
            batch_futures = [
                executor.submit(translate_single_text, index, text)
                for index, text in zip(batch_indices, batch_texts)
            ]

            # Collect results
            for future in concurrent.futures.as_completed(batch_futures):
                # Check cancellation
                if check_cancelled and check_cancelled():
                    logger.info("Translation cancelled")
                    break

                try:
                    index, translated_text = future.result()
                    translations[index] = translated_text
                    update_progress()
                except Exception as e:
                    logger.error(f"Batch processing error: {e}")

    # Add translation results to DataFrame
    df["translated_text"] = translations

    logger.info(f"Parallel translation complete: {len(df)} items")

    if progress_callback:
        try:
            progress_callback(0.8)  # Translation complete
        except Exception as e:
            logger.warning(f"Progress callback error: {e}")

    return df

def save_text_files(df: pd.DataFrame, output_dir: str, base_filename: str) -> tuple:
    """
    Save extracted and translated text to files.

    Args:
        df: DataFrame containing text data
        output_dir: Output directory
        base_filename: Base filename

    Returns:
        (extracted text file path, translated text file path)
    """
    logger.info(
        f"Saving text files: directory {output_dir}, base name {base_filename}"
    )

    # Create output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)

    # Generate filenames
    extracted_file = os.path.join(output_dir, f"{base_filename}_extracted.txt")
    translated_file = os.path.join(output_dir, f"{base_filename}_translated.txt")

    # Save extracted text
    with open(extracted_file, "w", encoding="utf-8") as f:
        for i, row in df.iterrows():
            # Format according to file type
            if "element_type" in row:
                if row["element_type"] == "paragraph":
                    f.write(f"=== Paragraph {row.get('para_id', i) + 1} ===\n")
                elif row["element_type"] == "table_cell":
                    f.write(
                        f"=== Table {row.get('table_id', 0) + 1}, Row {row.get('row_id', 0) + 1}, Cell {row.get('cell_id', 0) + 1} ===\n"
                    )
                elif row["element_type"] == "pdf_text":
                    f.write(
                        f"=== Page {row.get('page_num', 0) + 1}, Block {row.get('block_num', 0) + 1}, Line {row.get('line_num', 0) + 1} ===\n"
                    )
                elif row["element_type"] == "pdf_paragraph":
                    f.write(
                        f"=== Page {row.get('page_num', 0) + 1}, Paragraph {row.get('para_num', 0) + 1} ===\n"
                    )
            elif "slide_num" in row:
                f.write(
                    f"=== Slide {row['slide_num'] + 1}, Shape {row['shape_id'] + 1} ===\n"
                )
            else:
                f.write(f"=== Item {i + 1} ===\n")

            # Write original text
            f.write(f"{row['original_text']}\n\n")

            # Add hyperlink information if present
            if "hyperlink" in row and row["hyperlink"]:
                f.write(f"[Hyperlink: {row['hyperlink']}]\n\n")

    logger.info(f"Extracted text saved: {extracted_file}")

    # Save translated text (only if 'translated_text' column exists)
    if "translated_text" in df.columns:
        # Check if translated text is not empty
        if (
            not df["translated_text"].isna().all()
            and not (df["translated_text"] == "").all()
        ):
            with open(translated_file, "w", encoding="utf-8") as f:
                for i, row in df.iterrows():
                    # Format according to file type
                    if "element_type" in row:
                        if row["element_type"] == "paragraph":
                            f.write(f"=== Paragraph {row.get('para_id', i) + 1} ===\n")
                        elif row["element_type"] == "table_cell":
                            f.write(
                                f"=== Table {row.get('table_id', 0) + 1}, Row {row.get('row_id', 0) + 1}, Cell {row.get('cell_id', 0) + 1} ===\n"
                            )
                        elif row["element_type"] == "pdf_text":
                            f.write(
                                f"=== Page {row.get('page_num', 0) + 1}, Block {row.get('block_num', 0) + 1}, Line {row.get('line_num', 0) + 1} ===\n"
                            )
                        elif row["element_type"] == "pdf_paragraph":
                            f.write(
                                f"=== Page {row.get('page_num', 0) + 1}, Paragraph {row.get('para_num', 0) + 1} ===\n"
                            )
                    elif "slide_num" in row:
                        f.write(
                            f"=== Slide {row['slide_num'] + 1}, Shape {row['shape_id'] + 1} ===\n"
                        )
                    else:
                        f.write(f"=== Item {i + 1} ===\n")

                    # Use original text if translated text is None or empty
                    text_to_write = (
                        row["translated_text"]
                        if pd.notna(row["translated_text"]) and row["translated_text"]
                        else "[Not translated]"
                    )
                    f.write(f"{text_to_write}\n\n")

                    # Add hyperlink information if present
                    if "hyperlink" in row and row["hyperlink"]:
                        f.write(f"[Hyperlink: {row['hyperlink']}]\n\n")

            logger.info(f"Translated text saved: {translated_file}")
            return extracted_file, translated_file
        else:
            logger.warning(
                "Translated text is empty. Translated text file will not be saved."
            )
            return extracted_file, None
    else:
        logger.warning(
            "DataFrame does not have 'translated_text' column. Translated text will not be saved."
        )
        return extracted_file, None

#
# PPTX processing functions
#

def extract_text_from_pptx(
    pptx_path: str, progress_callback: Optional[Callable] = None
) -> pd.DataFrame:
    """
    Extract text from PPTX file and return as DataFrame.

    Args:
        pptx_path: PPTX file path
        progress_callback: Callback function to report progress

    Returns:
        DataFrame containing text data
    """
    logger.info(f"Extracting text from PPTX file '{pptx_path}'...")

    try:
        from pptx import Presentation

        prs = Presentation(pptx_path)
        text_data = []

        total_slides = len(prs.slides)
        logger.info(f"Number of slides: {total_slides}")

        # Save slide number, shape ID, and original text
        for slide_num, slide in enumerate(prs.slides):
            if progress_callback:
                try:
                    progress_callback(
                        slide_num / total_slides * 0.3
                    )  # Extraction is 30% of total
                except Exception as e:
                    logger.warning(f"Progress callback error: {e}")

            for shape_id, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    logger.debug(
                        f"Slide {slide_num+1}, Shape {shape_id+1}: {text[:50]}..."
                    )

                    # Get hyperlink information
                    hyperlink = None
                    if hasattr(shape, "click_action") and hasattr(
                        shape.click_action, "hyperlink"
                    ):
                        if (
                            hasattr(shape.click_action.hyperlink, "address")
                            and shape.click_action.hyperlink.address
                        ):
                            hyperlink = shape.click_action.hyperlink.address

                    text_data.append(
                        {
                            "slide_num": slide_num,
                            "shape_id": shape_id,
                            "original_text": text,
                            "hyperlink": hyperlink,
                        }
                    )

        # Convert to DataFrame
        df = pd.DataFrame(text_data)

        logger.info(f"Extraction complete: {len(df)} text elements found")

        if progress_callback:
            try:
                progress_callback(0.3)  # Extraction complete
            except Exception as e:
                logger.warning(f"Progress callback error: {e}")

        return df

    except ImportError:
        logger.error("python-pptx is not installed")
        raise ImportError(
            "python-pptx is not installed. Please install with pip install python-pptx."
        )
    except Exception as e:
        logger.error(f"Text extraction error from PPTX: {e}")
        raise

def reimport_text_to_pptx(
    pptx_path: str,
    df: pd.DataFrame,
    output_path: str,
    progress_callback: Optional[Callable] = None,
    check_cancelled: Optional[Callable] = None,
) -> None:
    """
    Re-import translated text into PPTX file.
    Maintain formatting and hyperlinks as much as possible.

    Args:
        pptx_path: Original PPTX file path
        df: DataFrame containing translated text
        output_path: Output PPTX file path
        progress_callback: Callback function to report progress
        check_cancelled: Function to check cancellation status
    """
    logger.info(
        f"Re-importing translated text into PPTX file: {output_path}"
    )

    try:
        from pptx import Presentation

        # Open presentation
        prs = Presentation(pptx_path)

        total_items = len(df)

        # Apply translated text to each slide and shape
        for i, (index, row) in enumerate(df.iterrows()):
            # Check cancellation
            if check_cancelled and check_cancelled():
                logger.info("PPTX re-import cancelled")
                return

            if progress_callback:
                try:
                    # Re-import is 80%~100% of total
                    progress = 0.8 + (i / total_items * 0.2)
                    progress_callback(progress)
                except Exception as e:
                    logger.warning(f"Progress callback error: {e}")

            slide_num = row["slide_num"]
            shape_id = row["shape_id"]
            translated_text = row["translated_text"]
            hyperlink = row.get("hyperlink")

            logger.debug(
                f"Applying translated text to Slide {slide_num+1}, Shape {shape_id+1}..."
            )

            # Identify corresponding slide and shape and replace text
            if slide_num < len(prs.slides):
                slide = prs.slides[slide_num]
                shapes = list(slide.shapes)
                if shape_id < len(shapes):
                    shape = shapes[shape_id]
                    if hasattr(shape, "text_frame"):
                        # Split translated text by paragraph
                        translated_paragraphs = translated_text.split("\n")
                        original_paragraphs = [p for p in shape.text_frame.paragraphs]

                        # Compare number of existing and translated paragraphs
                        if len(original_paragraphs) >= len(translated_paragraphs):
                            # Apply translated text to existing paragraphs (maintain formatting)
                            for p_idx, p_text in enumerate(translated_paragraphs):
                                if p_idx < len(original_paragraphs):
                                    # Replace text while maintaining original paragraph formatting
                                    original_p = original_paragraphs[p_idx]

                                    # Process runs in paragraph
                                    if len(original_p.runs) > 0:
                                        # Set text to first run
                                        original_p.runs[0].text = p_text

                                        # Empty remaining runs
                                        for run in original_p.runs[1:]:
                                            run.text = ""
                                    else:
                                        # Set paragraph text directly if no runs
                                        original_p.text = p_text

                            # Empty extra paragraphs
                            for p_idx in range(
                                len(translated_paragraphs), len(original_paragraphs)
                            ):
                                original_paragraphs[p_idx].text = ""
                        else:
                            # If more translated paragraphs
                            # Update existing paragraphs
                            for p_idx, original_p in enumerate(original_paragraphs):
                                if p_idx < len(translated_paragraphs):
                                    if len(original_p.runs) > 0:
                                        original_p.runs[0].text = translated_paragraphs[
                                            p_idx
                                        ]
                                        for run in original_p.runs[1:]:
                                            run.text = ""
                                    else:
                                        original_p.text = translated_paragraphs[p_idx]

                            # Add missing paragraphs
                            for p_idx in range(
                                len(original_paragraphs), len(translated_paragraphs)
                            ):
                                p = shape.text_frame.add_paragraph()
                                p.text = translated_paragraphs[p_idx]

                                # Copy style from first paragraph if possible
                                if original_paragraphs and hasattr(
                                    original_paragraphs[0], "alignment"
                                ):
                                    p.alignment = original_paragraphs[0].alignment

                        # Process hyperlinks
                        if hyperlink:
                            try:
                                # Due to python-pptx limitations, complete hyperlink reset may be difficult
                                # Set click action hyperlink if possible
                                if hasattr(shape, "click_action"):
                                    logger.debug(f"Maintaining hyperlink: {hyperlink}")
                                    shape.click_action.hyperlink.address = hyperlink
                            except Exception as e:
                                logger.warning(
                                    f"Failed to set hyperlink: {e}"
                                )

        # Save as new file
        logger.info(f"Saving translated PPTX file: {output_path}")
        prs.save(output_path)

        if progress_callback:
            try:
                progress_callback(1.0)  # Complete
            except Exception as e:
                logger.warning(f"Progress callback error: {e}")

    except Exception as e:
        logger.error(f"PPTX re-import error: {e}")
        raise

#
# DOCX processing functions
#

def extract_text_from_docx(
    docx_path: str, progress_callback: Optional[Callable] = None
) -> pd.DataFrame:
    """
    Extract text from Word (docx) file and return as DataFrame.

    Args:
        docx_path: Word file path
        progress_callback: Callback function to report progress

    Returns:
        DataFrame containing text data
    """
    logger.info(f"Extracting text from Word document '{docx_path}'...")

    try:
        from docx import Document

        doc = Document(docx_path)
        text_data = []

        # Calculate total number of paragraphs and tables (for progress display)
        total_elements = len(doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                total_elements += len(row.cells)

        logger.info(f"Number of elements: {total_elements}")

        # Extract text from paragraphs
        element_count = 0
        for para_id, para in enumerate(doc.paragraphs):
            if progress_callback:
                try:
                    progress_callback(
                        element_count / total_elements * 0.3
                    )  # Extraction is 30% of total
                except Exception as e:
                    logger.warning(f"Progress callback error: {e}")
                element_count += 1

            text = para.text.strip()
            if text:
                logger.debug(f"Paragraph {para_id+1}: {text[:50]}...")

                # Get style information
                style_name = para.style.name if para.style else "Normal"
                alignment = para.alignment if hasattr(para, "alignment") else None

                # Collect run formatting information in paragraph
                runs_info = []
                for run_id, run in enumerate(para.runs):
                    if run.text.strip():
                        run_info = {
                            "text": run.text,
                            "bold": run.bold,
                            "italic": run.italic,
                            "underline": run.underline,
                        }
                        runs_info.append(run_info)

                text_data.append(
                    {
                        "element_type": "paragraph",
                        "para_id": int(para_id),  # Ensure integer is used
                        "original_text": text,
                        "style_name": style_name,
                        "alignment": alignment,
                        "runs_info": runs_info,
                        "table_id": None,
                        "row_id": None,
                        "cell_id": None,
                    }
                )

        # Extract text from tables
        for table_id, table in enumerate(doc.tables):
            for row_id, row in enumerate(table.rows):
                for cell_id, cell in enumerate(row.cells):
                    if progress_callback:
                        try:
                            progress_callback(element_count / total_elements * 0.3)
                        except Exception as e:
                            logger.warning(f"Progress callback error: {e}")
                        element_count += 1

                    text = cell.text.strip()
                    if text:
                        logger.debug(
                            f"Table {table_id+1}, Row {row_id+1}, Cell {cell_id+1}: {text[:50]}..."
                        )

                        # Collect style information for paragraphs in cell
                        cell_paras_info = []
                        for p_id, p in enumerate(cell.paragraphs):
                            if p.text.strip():
                                para_info = {
                                    "text": p.text,
                                    "style_name": p.style.name if p.style else "Normal",
                                    "alignment": (
                                        p.alignment if hasattr(p, "alignment") else None
                                    ),
                                }
                                cell_paras_info.append(para_info)

                        text_data.append(
                            {
                                "element_type": "table_cell",
                                "para_id": None,
                                "original_text": text,
                                "style_name": None,
                                "alignment": None,
                                "runs_info": None,
                                "table_id": int(table_id),  # Ensure integer is used
                                "row_id": int(row_id),  # Ensure integer is used
                                "cell_id": int(cell_id),  # Ensure integer is used
                                "cell_paras_info": cell_paras_info,
                            }
                        )

        # Convert to DataFrame
        df = pd.DataFrame(text_data)

        logger.info(f"Extraction complete: {len(df)} text elements found")

        if progress_callback:
            try:
                progress_callback(0.3)  # Extraction complete
            except Exception as e:
                logger.warning(f"Progress callback error: {e}")

        return df

    except ImportError:
        logger.error("python-docx is not installed")
        raise ImportError(
            "python-docx is not installed. Please install with pip install python-docx."
        )
    except Exception as e:
        logger.error(f"Text extraction error from DOCX: {e}")
        raise

def reimport_text_to_docx(
    docx_path: str,
    df: pd.DataFrame,
    output_path: str,
    progress_callback: Optional[Callable] = None,
    check_cancelled: Optional[Callable] = None,
) -> None:
    """
    Re-import translated text into Word document.
    Maintain formatting as much as possible.

    Args:
        docx_path: Original Word file path
        df: DataFrame containing translated text
        output_path: Output Word file path
        progress_callback: Callback function to report progress
        check_cancelled: Function to check cancellation status
    """
    logger.info(f"Re-importing translated text into Word document: {output_path}")

    try:
        from docx import Document

        # Open document
        doc = Document(docx_path)

        total_items = len(df)

        # Apply translated text to paragraphs
        para_df = df[df["element_type"] == "paragraph"]
        for i, (_, row) in enumerate(para_df.iterrows()):
            # Check cancellation
            if check_cancelled and check_cancelled():
                logger.info("DOCX re-import cancelled")
                return

            if progress_callback:
                try:
                    # Re-import is 80%~100% of total
                    progress = 0.8 + (i / total_items * 0.2)
                    progress_callback(progress)
                except Exception as e:
                    logger.warning(f"Progress callback error: {e}")

            # Convert para_id to integer if it's a float
            try:
                para_id = int(row["para_id"])
            except (TypeError, ValueError):
                logger.warning(
                    f"Cannot convert paragraph ID '{row['para_id']}' to integer. Skipping."
                )
                continue

            translated_text = row["translated_text"]

            logger.debug(f"Applying translated text to Paragraph {para_id+1}...")

            # Identify corresponding paragraph and replace text
            if para_id < len(doc.paragraphs):
                para = doc.paragraphs[para_id]

                # Clear paragraph text
                for run in para.runs:
                    run.text = ""

                # Add translated text
                if para.runs:
                    # Set text to first run if existing runs
                    para.runs[0].text = translated_text
                else:
                    # Add new run if no runs
                    run = para.add_run(translated_text)

                    # Apply original formatting information if available
                    if (
                        "runs_info" in row
                        and row["runs_info"]
                        and len(row["runs_info"]) > 0
                    ):
                        run_info = row["runs_info"][0]
                        if run_info.get("bold"):
                            run.bold = True
                        if run_info.get("italic"):
                            run.italic = True
                        if run_info.get("underline"):
                            run.underline = True

        # Apply translated text to table cells
        table_df = df[df["element_type"] == "table_cell"]
        for i, (_, row) in enumerate(table_df.iterrows()):
            # Check cancellation
            if check_cancelled and check_cancelled():
                logger.info("DOCX re-import cancelled")
                return

            if progress_callback:
                try:
                    # Re-import is 80%~100% of total
                    progress = 0.8 + ((len(para_df) + i) / total_items * 0.2)
                    progress_callback(progress)
                except Exception as e:
                    logger.warning(f"Progress callback error: {e}")

            # Convert each ID to integer
            try:
                table_id = int(row["table_id"])
                row_id = int(row["row_id"])
                cell_id = int(row["cell_id"])
            except (TypeError, ValueError):
                logger.warning(
                    f"Cannot convert table ID '{row['table_id']}', row ID '{row['row_id']}', cell ID '{row['cell_id']}' to integer. Skipping."
                )
                continue

            translated_text = row["translated_text"]

            logger.debug(
                f"Applying translated text to Table {table_id+1}, Row {row_id+1}, Cell {cell_id+1}..."
            )

            # Identify corresponding cell and replace text
            try:
                if table_id < len(doc.tables):
                    table = doc.tables[table_id]
                    if row_id < len(table.rows):
                        row_obj = table.rows[row_id]
                        if cell_id < len(row_obj.cells):
                            cell = row_obj.cells[cell_id]

                            # Clear cell text
                            cell.text = ""

                            # Add translated text
                            # Create paragraphs based on cell paragraph information if available
                            if "cell_paras_info" in row and row["cell_paras_info"]:
                                paragraphs = translated_text.split("\n")
                                for p_idx, p_text in enumerate(paragraphs):
                                    if p_idx == 0:
                                        # First paragraph already exists
                                        p = cell.paragraphs[0]
                                    else:
                                        # Create additional paragraphs
                                        p = cell.add_paragraph()

                                    p.text = p_text

                                    # Apply style if original paragraph information available
                                    if p_idx < len(row["cell_paras_info"]):
                                        para_info = row["cell_paras_info"][p_idx]
                                        if para_info.get("alignment") is not None:
                                            p.alignment = para_info["alignment"]
                            else:
                                # Simply set text if no paragraph information
                                cell.text = translated_text
            except Exception as e:
                logger.warning(f"Error updating table cell: {e}")
                continue

        # Save as new file
        logger.info(f"Saving translated Word document: {output_path}")
        doc.save(output_path)

        if progress_callback:
            try:
                progress_callback(1.0)  # Complete
            except Exception as e:
                logger.warning(f"Progress callback error: {e}")

    except Exception as e:
        logger.error(f"DOCX re-import error: {e}")
        raise

#
# PDF processing functions (modified to match translator.py.old.py processing flow)
#

def convert_pdf_to_docx(pdf_path: str, docx_path: str) -> str:
    """
    Convert PDF file to DOCX.
    Prioritize pdf2docx library, use LibreOffice on failure.
    
    Args:
        pdf_path: Input PDF file path
        docx_path: Output DOCX file path
        
    Returns:
        Converted DOCX file path
    """
    logger.info(f"Converting PDF to DOCX: {pdf_path} -> {docx_path}")
    
    # Method 1: Use pdf2docx library
    try:
        from pdf2docx import Converter
        
        logger.info("Converting PDF to DOCX using pdf2docx...")
        
        # Convert PDF to DOCX
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()
        
        # Verify conversion result
        if os.path.exists(docx_path) and os.path.getsize(docx_path) > 0:
            logger.info(f"PDF→DOCX conversion complete using pdf2docx: {docx_path}")
            return docx_path
        else:
            logger.warning("DOCX not created by pdf2docx")
            raise Exception("DOCX not created")
            
    except ImportError:
        logger.warning("pdf2docx not installed. Trying LibreOffice.")
    except Exception as e:
        logger.warning(f"PDF→DOCX conversion failed using pdf2docx: {e}. Trying LibreOffice.")
    
    # Method 2: Convert PDF to DOCX using LibreOffice
    try:
        logger.info("Converting PDF to DOCX using LibreOffice...")
        
        # Search for LibreOffice path
        libreoffice_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS
            "/usr/local/bin/soffice",
            "/usr/bin/soffice",
            "libreoffice",
            "soffice"
        ]
        
        libreoffice_cmd = None
        for path in libreoffice_paths:
            try:
                if os.path.exists(path) or path in ["libreoffice", "soffice"]:
                    # Verify command existence
                    result = subprocess.run([path, "--version"], 
                                          capture_output=True, text=True, timeout=10)
                    if result.returncode == 0:
                        libreoffice_cmd = path
                        break
            except (subprocess.TimeoutExpired, FileNotFoundError, subprocess.CalledProcessError):
                continue
        
        if not libreoffice_cmd:
            raise FileNotFoundError("LibreOffice not found")
        
        # Convert PDF to DOCX using LibreOffice
        output_dir = os.path.dirname(docx_path)
        cmd = [
            libreoffice_cmd,
            "--headless", "--convert-to", "docx",
            "--outdir", output_dir, 
            pdf_path
        ]
        
        result = subprocess.run(cmd, check=True, timeout=120, 
                              capture_output=True, text=True)
        
        # LibreOffice uses original filename, rename if necessary
        generated_docx = os.path.join(
            output_dir,
            os.path.splitext(os.path.basename(pdf_path))[0] + ".docx"
        )
        
        if generated_docx != docx_path and os.path.exists(generated_docx):
            shutil.move(generated_docx, docx_path)
        
        if os.path.exists(docx_path) and os.path.getsize(docx_path) > 0:
            logger.info(f"PDF→DOCX conversion complete using LibreOffice: {docx_path}")
            return docx_path
        else:
            logger.warning("DOCX not created by LibreOffice")
            raise Exception("DOCX not created")
            
    except Exception as e:
        logger.warning(f"PDF→DOCX conversion failed using LibreOffice: {e}")
        
        # Method 3: Extract text directly from PDF and create DOCX
        try:
            logger.info("Extracting text directly from PDF and creating DOCX...")
            
            # Extract text from PDF
            text_content = []
            
            # Try pdfplumber
            try:
                import pdfplumber
                with pdfplumber.open(pdf_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text and text.strip():
                            text_content.append(f"=== Page {page_num + 1} ===\n{text}\n")
            except Exception as e:
                logger.warning(f"Extraction failed with pdfplumber: {e}")
                
                # Try PyPDF2
                try:
                    import PyPDF2
                    with open(pdf_path, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        for page_num, page in enumerate(pdf_reader.pages):
                            text = page.extract_text()
                            if text and text.strip():
                                text_content.append(f"=== Page {page_num + 1} ===\n{text}\n")
                except Exception as e2:
                    logger.error(f"Extraction also failed with PyPDF2: {e2}")
                    raise ValueError(f"Failed to extract text from PDF: {str(e2)}")
            
            if not text_content:
                raise ValueError("Could not extract text from PDF")
            
            # Create DOCX file
            from docx import Document
            doc = Document()
            
            for content in text_content:
                paragraphs = content.split('\n')
                for paragraph_text in paragraphs:
                    if paragraph_text.strip():
                        doc.add_paragraph(paragraph_text)
            
            doc.save(docx_path)
            
            if os.path.exists(docx_path) and os.path.getsize(docx_path) > 0:
                logger.info(f"DOCX creation complete using text extraction: {docx_path}")
                return docx_path
            else:
                raise ValueError("Failed to create DOCX file")
                
        except Exception as e2:
            logger.error(f"DOCX creation also failed using text extraction: {e2}")
            raise ValueError(f"Failed to convert PDF to DOCX: {str(e2)}")

def convert_docx_to_pdf(docx_path: str, pdf_path: str) -> str:
    """
    Convert DOCX file to PDF.
    Try multiple methods with fallback processing.
    
    Args:
        docx_path: Input DOCX file path
        pdf_path: Output PDF file path
        
    Returns:
        Converted PDF file path
    """
    import sys
    import threading
    import time
    
    logger.info(f"Converting DOCX to PDF: {docx_path} -> {pdf_path}")
    
    # Method 1: Try docx2pdf
    try:
        from docx2pdf import convert
        
        logger.info("Converting DOCX to PDF using docx2pdf...")
        
        def convert_with_timeout():
            try:
                convert(docx_path, pdf_path)
                return True
            except Exception as e:
                logger.error(f"docx2pdf conversion error: {e}")
                return False
        
        # Execute conversion in separate thread (with timeout)
        thread = threading.Thread(target=convert_with_timeout)
        thread.daemon = True
        thread.start()
        
        # Wait up to 60 seconds
        timeout = 60
        start_time = time.time()
        while thread.is_alive() and time.time() - start_time < timeout:
            time.sleep(1)
        
        if thread.is_alive():
            logger.warning(f"docx2pdf timed out ({timeout} seconds)")
            raise TimeoutError(f"docx2pdf timed out ({timeout} seconds)")
        
        # Verify conversion result
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            logger.info(f"DOCX conversion complete using docx2pdf: {pdf_path}")
            return pdf_path
        else:
            logger.warning("PDF not created by docx2pdf")
            raise Exception("PDF not created")
            
    except ImportError:
        logger.warning("docx2pdf not installed. Trying alternative method.")
    except Exception as e:
        logger.warning(f"Conversion failed using docx2pdf: {e}. Trying alternative method.")
    
    # Method 2: Try LibreOffice
    try:
        logger.info("Converting DOCX to PDF using LibreOffice...")
        
        # Search for LibreOffice path
        libreoffice_paths = [
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS
            "/usr/local/bin/soffice",
            "/usr/bin/soffice",
            "libreoffice",
            "soffice"
        ]
        
        libreoffice_cmd = None
        for path in libreoffice_paths:
            try:
                if os.path.exists(path) or path in ["libreoffice", "soffice"]:
                    # Verify command existence
                    result = subprocess.run([path, "--version"], 
                                          capture_output=True, text=True, timeout=10)
                    if result.returncode == 0:
                        libreoffice_cmd = path
                        break
            except (subprocess.TimeoutExpired, FileNotFoundError, subprocess.CalledProcessError):
                continue
        
        if not libreoffice_cmd:
            raise FileNotFoundError("LibreOffice not found")
        
        # Convert DOCX to PDF using LibreOffice
        cmd = [
            libreoffice_cmd,
            "--headless", "--convert-to", "pdf",
            "--outdir", os.path.dirname(pdf_path), 
            docx_path
        ]
        
        result = subprocess.run(cmd, check=True, timeout=120, 
                              capture_output=True, text=True)
        
        # LibreOffice uses original filename, rename if necessary
        generated_pdf = os.path.join(
            os.path.dirname(pdf_path),
            os.path.splitext(os.path.basename(docx_path))[0] + ".pdf"
        )
        
        if generated_pdf != pdf_path and os.path.exists(generated_pdf):
            shutil.move(generated_pdf, pdf_path)
        
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            logger.info(f"DOCX conversion complete using LibreOffice: {pdf_path}")
            return pdf_path
        else:
            logger.warning("PDF not created by LibreOffice")
            raise Exception("PDF not created")
            
    except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired) as e:
        logger.warning(f"Conversion failed using LibreOffice: {e}")
    
    # Method 3: Simple PDF generation using reportlab
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from docx import Document
        
        logger.info("Generating simple PDF using reportlab...")
        
        # Extract text from DOCX
        doc = Document(docx_path)
        text_content = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        # Generate PDF
        pdf_doc = SimpleDocTemplate(pdf_path, pagesize=letter)
        styles = getSampleStyleSheet()
        flowables = []
        
        for text in text_content:
            try:
                # Remove problematic characters
                safe_text = ''.join(c for c in text if ord(c) < 65536)
                p = Paragraph(safe_text, styles["Normal"])
                flowables.append(p)
                flowables.append(Spacer(1, 12))
            except Exception as e:
                logger.warning(f"Failed to create paragraph: {e}")
                continue
        
        pdf_doc.build(flowables)
        
        if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
            logger.info(f"Simple PDF generation complete using reportlab: {pdf_path}")
            logger.warning("Note: Simple PDF has limited formatting")
            return pdf_path
        else:
            logger.warning("PDF not created by reportlab")
            raise Exception("PDF not created")
            
    except ImportError:
        logger.warning("reportlab not installed")
    except Exception as e:
        logger.warning(f"Conversion failed using reportlab: {e}")
    
    # All methods failed
    logger.error("All PDF conversion methods failed")
    
    # Last resort: Copy DOCX file as is
    docx_copy_path = pdf_path.replace(".pdf", ".docx")
    shutil.copy(docx_path, docx_copy_path)
    logger.warning(f"PDF conversion failed, copied DOCX file: {docx_copy_path}")
    
    # Create text file with error message
    error_txt_path = pdf_path.replace(".pdf", "_error.txt")
    with open(error_txt_path, "w", encoding="utf-8") as f:
        f.write(f"PDF conversion failed. DOCX file saved instead: {docx_copy_path}\n")
        f.write("All conversion methods (docx2pdf, LibreOffice, reportlab) failed.\n")
    
    return docx_copy_path

def translate_pdf(
    input_path: str,
    output_path: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    progress_callback: Optional[Callable] = None,
    save_text_files_flag: bool = True,
    ai_instruction: str = "",
    check_cancelled: Optional[Callable] = None,
) -> tuple:
    """
    Translate PDF file.
    Follow translator.py.old.py processing flow:
    1. Convert PDF to DOCX with convert_pdf_to_docx()
    2. Execute DOCX translation with translate_docx()
    3. Convert translated DOCX to PDF with convert_docx_to_pdf()
    
    Args:
        input_path: Input PDF file path
        output_path: Output PDF file path
        api_key: GenAI Hub API key
        model: Model name to use
        source_lang: Source language code
        target_lang: Target language code
        progress_callback: Callback function to report progress
        save_text_files_flag: Whether to save text files
        ai_instruction: Supplementary instructions for AI
        check_cancelled: Function to check cancellation status
        
    Returns:
        (extracted text file path, translated text file path)
    """
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()
    
    logger.info(f"Starting PDF translation: {input_path} -> {output_path}")
    
    if ai_instruction:
        logger.info(f"Supplementary instructions for AI: {ai_instruction}")
    
    # Generate temporary file paths
    temp_dir = os.path.dirname(output_path)
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    temp_docx_input = os.path.join(temp_dir, f"{base_name}_temp_input.docx")
    temp_docx_output = os.path.join(temp_dir, f"{base_name}_temp_output.docx")
    
    try:
        # Check cancellation
        if check_cancelled and check_cancelled():
            logger.info("Translation cancelled")
            return None, None
        
        # Update progress
        if progress_callback:
            progress_callback(0.05)
            
        # 1. Convert PDF to DOCX
        logger.info("Converting PDF to DOCX...")
        convert_pdf_to_docx(input_path, temp_docx_input)
        
        # Check cancellation
        if check_cancelled and check_cancelled():
            logger.info("Translation cancelled")
            return None, None
        
        if progress_callback:
            progress_callback(0.15)
        
        # 2. Translate DOCX
        logger.info("Translating DOCX...")
        extracted_file, translated_file = translate_docx(
            temp_docx_input,
            temp_docx_output,
            api_key,
            model,
            source_lang,
            target_lang,
            # Adjust progress callback to 15%~85% range
            lambda p: progress_callback(0.15 + p * 0.7) if progress_callback else None,
            save_text_files_flag,
            ai_instruction,
            check_cancelled
        )
        
        # Check cancellation
        if check_cancelled and check_cancelled():
            logger.info("Translation cancelled")
            return extracted_file, translated_file
        
        # 3. Convert translated DOCX to PDF
        logger.info("Converting translated DOCX to PDF...")
        final_output = convert_docx_to_pdf(temp_docx_output, output_path)
        
        if progress_callback:
            progress_callback(1.0)
            
        # If final output is DOCX file (PDF conversion failed)
        if final_output.endswith('.docx'):
            logger.warning("PDF conversion failed. DOCX file will be provided.")
            # Move DOCX file to appropriate location
            if final_output != output_path.replace('.pdf', '.docx'):
                shutil.move(final_output, output_path.replace('.pdf', '.docx'))
        
        logger.info(f"PDF translation complete: {output_path}")
        return extracted_file, translated_file
        
    except Exception as e:
        logger.error(f"PDF translation error: {e}")
        raise
    finally:
        # Delete temporary files
        try:
            if os.path.exists(temp_docx_input):
                os.remove(temp_docx_input)
                logger.debug(f"Deleted temporary file: {temp_docx_input}")
            if os.path.exists(temp_docx_output):
                os.remove(temp_docx_output)
                logger.debug(f"Deleted temporary file: {temp_docx_output}")
        except Exception as e:
            logger.warning(f"Failed to delete temporary files: {e}")

def translate_xlsx(
    input_path: str,
    output_path: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    progress_callback: Optional[Callable] = None,
    save_text_files_flag: bool = True,
    ai_instruction: str = "",
    check_cancelled: Optional[Callable] = None,
) -> tuple:
    """
    Translate XLSX file.

    Args:
        input_path: Input Excel file path
        output_path: Output Excel file path
        api_key: GenAI Hub API key
        model: Model name to use
        source_lang: Source language code
        target_lang: Target language code
        progress_callback: Callback function to report progress
        save_text_files_flag: Whether to save text files
        ai_instruction: Supplementary instructions for AI
        check_cancelled: Function to check cancellation status

    Returns:
        (extracted text file path, translated text file path)
    """
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()
    
    logger.info(f"Starting Excel translation: {input_path} -> {output_path}")

    try:
        # Display progress
        if progress_callback:
            progress_callback(0.1)

        # Read Excel file
        excel = pd.read_excel(input_path, sheet_name=None, engine="openpyxl")
        total_sheets = len(excel)
        
        # Extract text data and convert to DataFrame
        text_data = []
        for sheet_idx, (sheet_name, df) in enumerate(excel.items()):
            if check_cancelled and check_cancelled():
                return None, None
                
            for row_idx, row in df.iterrows():
                for col_idx, value in row.items():
                    if isinstance(value, str) and value.strip():
                        text_data.append({
                            'sheet_name': sheet_name,
                            'row': row_idx,
                            'column': col_idx,
                            'original_text': value.strip()
                        })
            
            if progress_callback:
                progress_callback(0.1 + (sheet_idx / total_sheets * 0.2))

        # Convert text data to DataFrame
        text_df = pd.DataFrame(text_data)

        # Save text files
        output_dir = os.path.dirname(output_path)
        base_filename = os.path.splitext(os.path.basename(output_path))[0]
        extracted_file = None
        translated_file = None

        if save_text_files_flag:
            extracted_file, _ = save_text_files(text_df, output_dir, base_filename)

        # Translate text
        text_df = translate_dataframe(
            text_df,
            api_key,
            model,
            source_lang,
            target_lang,
            lambda p: progress_callback(0.3 + p * 0.5) if progress_callback else None,
            ai_instruction,
            check_cancelled
        )

        if save_text_files_flag and 'translated_text' in text_df.columns:
            _, translated_file = save_text_files(text_df, output_dir, base_filename)

        # Reflect translation results in Excel
        translated_excel = excel.copy()
        for _, row in text_df.iterrows():
            if check_cancelled and check_cancelled():
                return extracted_file, translated_file
                
            if 'translated_text' in row and pd.notna(row['translated_text']):
                sheet = translated_excel[row['sheet_name']]
                sheet.at[row['row'], row['column']] = row['translated_text']

        # Save translated Excel file
        with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
            for sheet_name, df in translated_excel.items():
                df.to_excel(writer, sheet_name=sheet_name, index=False)

        if progress_callback:
            progress_callback(1.0)

        return extracted_file, translated_file

    except Exception as e:
        logger.error(f"Excel translation error: {e}")
        raise

def translate_pptx(
    input_path: str,
    output_path: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    progress_callback: Optional[Callable] = None,
    save_text_files_flag: bool = True,
    ai_instruction: str = "",
    check_cancelled: Optional[Callable] = None,
) -> tuple:
    """
    Translate PPTX file.

    Args:
        input_path: Input PPTX file path
        output_path: Output PPTX file path
        api_key: GenAI Hub API key
        model: Model name to use
        source_lang: Source language code
        target_lang: Target language code
        progress_callback: Callback function to report progress
        save_text_files_flag: Whether to save text files
        ai_instruction: Supplementary instructions for AI
        check_cancelled: Function to check cancellation status

    Returns:
        (extracted text file path, translated text file path)
    """
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()
    
    logger.info(f"Starting PPTX translation: {input_path} -> {output_path}")
    logger.info(f"Languages: {source_lang} -> {target_lang}, Model: {model}")

    if ai_instruction:
        logger.info(f"Supplementary instructions for AI: {ai_instruction}")

    # Get API key from environment variable
    if not api_key:
        api_key = os.environ.get("GENAI_HUB_API_KEY")

    if not api_key:
        raise ValueError(
            "API key is not set. Please check environment variable GENAI_HUB_API_KEY."
        )

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return None, None

    # 1. Text extraction
    df = extract_text_from_pptx(input_path, progress_callback)

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return None, None

    # Get output directory and base filename
    output_dir = os.path.dirname(output_path)
    base_filename = os.path.splitext(os.path.basename(output_path))[0]

    # Save extracted text
    extracted_file = None
    translated_file = None

    if save_text_files_flag:
        extracted_file, _ = save_text_files(df, output_dir, base_filename)

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, None

    # 2. Text translation
    df = translate_dataframe(
        df,
        api_key,
        model,
        source_lang,
        target_lang,
        progress_callback,
        ai_instruction,
        check_cancelled,
    )

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, None

    # Save translated text
    if save_text_files_flag and "translated_text" in df.columns:
        _, translated_file = save_text_files(df, output_dir, base_filename)

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, translated_file

    # 3. Re-import translated text
    reimport_text_to_pptx(
        input_path, df, output_path, progress_callback, check_cancelled
    )

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, translated_file

    logger.info(f"PPTX translation complete: {output_path}")

    return extracted_file, translated_file

def translate_docx(
    input_path: str,
    output_path: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    progress_callback: Optional[Callable] = None,
    save_text_files_flag: bool = True,
    ai_instruction: str = "",
    check_cancelled: Optional[Callable] = None,
) -> tuple:
    """
    Translate Word (docx) file.

    Args:
        input_path: Input Word file path
        output_path: Output Word file path
        api_key: GenAI Hub API key
        model: Model name to use
        source_lang: Source language code
        target_lang: Target language code
        progress_callback: Callback function to report progress
        save_text_files_flag: Whether to save text files
        ai_instruction: Supplementary instructions for AI
        check_cancelled: Function to check cancellation status

    Returns:
        (extracted text file path, translated text file path)
    """
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()
    
    logger.info(f"Starting Word document translation: {input_path} -> {output_path}")
    logger.info(f"Languages: {source_lang} -> {target_lang}, Model: {model}")

    if ai_instruction:
        logger.info(f"Supplementary instructions for AI: {ai_instruction}")

    # Get API key from environment variable
    if not api_key:
        api_key = os.environ.get("GENAI_HUB_API_KEY")

    if not api_key:
        raise ValueError(
            "API key is not set. Please check environment variable GENAI_HUB_API_KEY."
        )

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return None, None

    # 1. Text extraction
    df = extract_text_from_docx(input_path, progress_callback)

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return None, None

    # Get output directory and base filename
    output_dir = os.path.dirname(output_path)
    base_filename = os.path.splitext(os.path.basename(output_path))[0]

    # Save extracted text
    extracted_file = None
    translated_file = None

    if save_text_files_flag:
        extracted_file, _ = save_text_files(df, output_dir, base_filename)

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, None

    # 2. Text translation
    df = translate_dataframe(
        df,
        api_key,
        model,
        source_lang,
        target_lang,
        progress_callback,
        ai_instruction,
        check_cancelled,
    )

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, None

    # Save translated text
    if save_text_files_flag and "translated_text" in df.columns:
        _, translated_file = save_text_files(df, output_dir, base_filename)

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, translated_file

    # 3. Re-import translated text
    reimport_text_to_docx(
        input_path, df, output_path, progress_callback, check_cancelled
    )

    # Check cancellation
    if check_cancelled and check_cancelled():
        logger.info("Translation cancelled")
        return extracted_file, translated_file

    logger.info(f"Word document translation complete: {output_path}")

    return extracted_file, translated_file

def check_libreoffice() -> Tuple[bool, str]:
    """
    Check LibreOffice availability.

    Returns:
        (availability, version information or error message)
    """
    try:
        result = subprocess.run(
            ["libreoffice", "--version"], capture_output=True, text=True
        )

        if result.returncode == 0:
            version = result.stdout.strip()
            logger.info(f"LibreOffice version: {version}")
            return True, version
        else:
            error_msg = f"LibreOffice command returned error: {result.stderr}"
            logger.error(error_msg)
            return False, error_msg

    except FileNotFoundError:
        error_msg = "LibreOffice command not found. May not be installed."
        logger.error(error_msg)
        return False, error_msg
    except Exception as e:
        error_msg = f"LibreOffice check error: {e}"
        logger.error(error_msg)
        return False, error_msg

def detect_file_type(file_path: str) -> str:
    """
    Detect file format from file extension.

    Args:
        file_path: File path

    Returns:
        File format ('pptx', 'docx', 'pdf')
    """
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        return "pptx"
    elif ext == ".docx":
        return "docx"
    elif ext == ".pdf":
        return "pdf"
    elif ext == ".xlsx":
        return "xlsx"
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def translate_document(
    input_path: str,
    output_path: str,
    api_key: str = None,
    model: str = None,
    source_lang: str = "en",
    target_lang: str = "ja",
    progress_callback: Optional[Callable] = None,
    save_text_files_flag: bool = True,
    ai_instruction: str = "",
    check_cancelled: Optional[Callable] = None,
) -> tuple:
    """
    Translate document (PPTX, DOCX, PDF, XLSX).
    """
    # Auto-select model if not specified
    if model is None:
        from app.config import get_default_model
        model = get_default_model()
    
    logger.info(f"Starting document translation: {input_path} -> {output_path}")

    # Load cache at startup
    if TRANSLATION_CONFIG["use_cache"] and not TRANSLATION_CACHE:
        load_translation_cache()

    file_type = detect_file_type(input_path)

    try:
        # Check cancellation at each processing step
        if check_cancelled and check_cancelled():
            logger.info("Translation cancelled by user")
            return None, None

        # Process according to file format
        if file_type == "xlsx":
            result = translate_xlsx(
                input_path,
                output_path,
                api_key,
                model,
                source_lang,
                target_lang,
                progress_callback,
                save_text_files_flag,
                ai_instruction,
                check_cancelled
            )
        elif file_type == "pptx":
            result = translate_pptx(
                input_path,
                output_path,
                api_key,
                model,
                source_lang,
                target_lang,
                progress_callback,
                save_text_files_flag,
                ai_instruction,
                check_cancelled
            )
        elif file_type == "docx":
            result = translate_docx(
                input_path,
                output_path,
                api_key,
                model,
                source_lang,
                target_lang,
                progress_callback,
                save_text_files_flag,
                ai_instruction,
                check_cancelled
            )
        elif file_type == "pdf":
            result = translate_pdf(
                input_path,
                output_path,
                api_key,
                model,
                source_lang,
                target_lang,
                progress_callback,
                save_text_files_flag,
                ai_instruction,
                check_cancelled
            )
        else:
            raise ValueError(f"Unsupported file format: {file_type}")

        # Save cache after translation complete
        if TRANSLATION_CONFIG["use_cache"]:
            save_translation_cache()

        return result

    except Exception as e:
        logger.error(f"Translation error: {e}")
        if TRANSLATION_CONFIG["use_cache"]:
            save_translation_cache()
        raise

# Test code
if __name__ == "__main__":
    # Get API key from environment variable
    api_key = os.environ.get("GENAI_HUB_API_KEY")

    if not api_key:
        print("Environment variable GENAI_HUB_API_KEY is not set.")
        api_key = input("Enter GenAI Hub API key: ")

    # Simple test translation
    test_text = "Hello, world. This is a translation test."
    print(f"Test translation: {test_text}")

    translated = translate_text_with_cache(
        test_text, api_key, source_lang="en", target_lang="ja"
    )

    print(f"Translation result: {translated}")

# Export functions and variables
__all__ = [
    "get_available_models",
    "fetch_available_models",
    "LANGUAGES",
    "TRANSLATION_CONFIG",
    "translate_pptx",
    "translate_docx",
    "translate_pdf",
    "translate_xlsx",
    "translate_document",
    "load_translation_cache",
    "save_translation_cache",
]
